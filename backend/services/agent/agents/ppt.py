import base64
import logging
import os
import random
import re
import textwrap
import time
from concurrent.futures import ThreadPoolExecutor
from io import BytesIO
from typing import Any, Dict, List, Literal, Optional, Tuple

import requests
from PIL import Image, ImageDraw, ImageFilter, ImageOps

from configs.llm import get_llm
from configs.models import PPT_MODEL
from langchain_core.messages import HumanMessage, SystemMessage
from pydantic import BaseModel, Field

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger("ppt_agent")
logger.setLevel(logging.INFO)
if not logging.getLogger().handlers:
    logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(name)s: %(message)s")


# ---------------------------------------------------------------------------
# Schemas (unchanged -- kept identical so this is a drop-in replacement for
# whatever already calls get_ppt_chain() / ppt_agent())
# ---------------------------------------------------------------------------

class ThemeConfig(BaseModel):
    bg_color: str = Field(description="Background color hex code (e.g., 0F172A). DO NOT INCLUDE #.")
    text_color: str = Field(description="Main text color hex code. DO NOT INCLUDE #.")
    accent_1: str = Field(description="Primary accent color hex. DO NOT INCLUDE #.")
    accent_2: str = Field(description="Secondary accent color hex. DO NOT INCLUDE #.")
    panel_color: str = Field(description="Panel/Card background color hex. DO NOT INCLUDE #.")
    font_name: str = Field(default="Arial", description="Primary font name")
    style_name: str = Field(description="Name of the style, e.g., Cyberpunk, Corporate Minimal, etc.")


class VisualElement(BaseModel):
    visual_type: Literal["image", "icon", "architecture", "timeline", "process", "dashboard", "none"] = Field(
        description="Type of visual to include. Use 'architecture' for tech stacks, 'process' for workflows, 'dashboard' for metrics, 'image' for photo visuals."
    )
    description: str = Field(description="Describe the visual, e.g., 'A 3-step deployment pipeline' or 'A photo of a modern office'")
    search_keyword: str = Field(description="A clean 1-2 word English keyword for fetching stock photos, e.g. 'technology', 'city', 'office', 'people', 'data'.")


class DataItem(BaseModel):
    title: str = Field(description="Title of the card, step, metric, or component. Keep short (2-4 words).")
    description: str = Field(description="Short concise description (10-18 words max). DO NOT write long paragraphs.")


class SlideContent(BaseModel):
    layout: Literal[
        "title",
        "hero",
        "two_column",
        "cards",
        "timeline",
        "process_flow",
        "comparison",
        "dashboard",
        "architecture",
        "roadmap",
        "infographic",
        "bullets"
    ] = Field(description="The layout template to use for this slide. NEVER use bullets if a visual layout (cards, timeline, architecture, etc.) works better.")
    title: str = Field(description="Short slide title (4-7 words max).")
    subtitle: str = Field(default="", description="Optional supporting line (6-10 words max).")
    bullets: List[str] = Field(default_factory=list, description="List of concise bullet points (max 4 bullets, 8-15 words per bullet).")
    data_items: List[DataItem] = Field(
        default_factory=list,
        description="Key-value pairs for cards, comparison, dashboard or timeline. Maximum 4 items."
    )
    visual: VisualElement
    speaker_notes: str = ""


class PresentationDeck(BaseModel):
    title: str
    audience: str
    theme: ThemeConfig
    slides: List[SlideContent]


SYSTEM_PROMPT = """
You are the elite AI Presentation Architect for OmniMindAI.
Your job is to analyze the user's request and design a stunning, highly visual, publication-ready PowerPoint deck.

CRITICAL FORMATTING & LAYOUT RULES:
1. NO BORING TEMPLATES: Generate a unique `ThemeConfig` based on the prompt's topic. For a cyber security presentation, use dark backgrounds with neon accents. For healthcare, use clean whites and calming blues. DO NOT include '#' in hex codes.
2. NO WALLS OF TEXT: Keep slide text extremely concise and punchy.
   - Slide Title: 4-7 words max.
   - Slide Subtitle: 6-10 words max.
   - Card Item Title: 2-4 words.
   - Card Item Description: 10-18 words max. Crisp bullet-style summaries ONLY.
3. VISUAL INTELLIGENCE & IMAGES:
   - Always provide a valid `VisualElement`.
   - Provide a simple 1-2 word English keyword in `search_keyword` (e.g., "technology", "people", "business", "city", "growth", "code").
   - Use `visual_type="image"` for slides where photos fit well.
4. SLIDE LAYOUTS:
   - `title`: Opening title slide.
   - `hero`: Bold statement with featured image.
   - `two_column`: Content on left, image or structure on right.
   - `cards` / `infographic`: 3-4 distinct cards.
   - `timeline` / `roadmap`: Sequence steps.
   - `process_flow` / `architecture`: Tech stack or workflow steps.
   - `comparison`: Compare 2 items.
   - `dashboard`: Display 4 metrics/KPIs.
5. DASHBOARD METRIC CARDS: for `layout="dashboard"`, each `DataItem` is one KPI card.
   - Put the short headline stat (a number, percentage, ranking, or count -- e.g. "65%+", "120+", "#3", "1.4B+") in `title`.
   - Put the one-sentence explanation of what that stat means in `description`.
   - Keep the stat itself under ~10 characters so it reads as a bold headline number, not a phrase.

Respond perfectly in JSON format following the schema.
"""


def get_ppt_chain():
    ppt_llm = get_llm(
        model=PPT_MODEL,
        temperature=0.35,
        streaming=False,
    )
    return ppt_llm.with_structured_output(PresentationDeck)


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

IMAGE_TIMEOUT = 8          # seconds per HTTP attempt
IMAGE_MAX_RETRIES = 2      # attempts per source before moving to the next one
IMAGE_MAX_WORKERS = 6      # parallel image downloads during prefetch
IMAGE_DPI = 150            # render resolution used when cropping fetched/placeholder photos
IMAGE_CACHE_MAX = 300      # crude cap so a long-lived server process doesn't grow this forever

# Optional API keys. If present in the environment, these sources are tried
# FIRST because they return real, keyword-relevant photos. If absent, the
# agent still works end-to-end: it silently falls back to keyless sources
# and finally to a locally generated placeholder, so a slide never ships
# with a hole in it either way.
PEXELS_API_KEY = os.environ.get("PEXELS_API_KEY", "").strip()
PIXABAY_API_KEY = os.environ.get("PIXABAY_API_KEY", "").strip()
UNSPLASH_ACCESS_KEY = os.environ.get("UNSPLASH_ACCESS_KEY", "").strip()

HTTP_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                  "(KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36",
    "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
}

_HEX_RE = re.compile(r'^[0-9A-Fa-f]{6}$')
_STAT_RE = re.compile(r'\d')


# ---------------------------------------------------------------------------
# Text helpers: truncation + dynamic font fitting
# ---------------------------------------------------------------------------

def truncate_chars(text: Optional[str], max_chars: int) -> str:
    """Hard ceiling on a string's length, breaking on a word boundary where
    possible. This is defense-in-depth: fit_font_size() adapts to whatever
    text it's given, but a truly pathological string (a model ignoring the
    prompt's word limits) should never reach the renderer un-capped."""
    text = (text or "").strip()
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars].rsplit(" ", 1)[0]
    return (cut or text[:max_chars]).rstrip(",.;:-") + "…"


def fit_font_size(
    text: str,
    width_in: float,
    height_in: float,
    max_pt: int = 16,
    min_pt: int = 9,
    bold: bool = False,
    line_spacing: float = 1.24,
    avg_char_width_ratio: float = 0.52,
    h_margin_in: float = 0.36,
    v_margin_in: float = 0.36,
) -> int:
    """Pick the largest font size (points) at which `text` still word-wraps
    to fit inside a box of the given size.

    This is a deliberately simple simulation -- estimated chars-per-line
    from an average glyph-width ratio, rather than exact font metrics --
    because we can't guarantee which font is installed wherever the deck is
    opened. The ratio is tuned conservatively (errs toward a slightly
    smaller font) so the estimate stays safely inside the box rather than
    right at the edge of it.
    """
    text = (text or "").strip()
    if not text:
        return max_pt

    bold_factor = 1.06 if bold else 1.0
    usable_width = max(width_in - h_margin_in, 0.4)
    usable_height = max(height_in - v_margin_in, 0.22)

    for size in range(max_pt, min_pt - 1, -1):
        avg_char_w_in = (size * avg_char_width_ratio * bold_factor) / 72.0
        chars_per_line = max(1, int(usable_width / avg_char_w_in))
        num_lines = len(textwrap.wrap(text, width=chars_per_line)) or 1
        line_height_in = (size * line_spacing) / 72.0
        if num_lines * line_height_in <= usable_height:
            return size

    return min_pt


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def hex_to_tuple(hex_str: Optional[str], fallback: str = "6B7280") -> Tuple[int, int, int]:
    cleaned = (hex_str or "").strip().lstrip('#')
    if not _HEX_RE.match(cleaned):
        cleaned = fallback
    return (int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))


def hex_to_rgb(hex_str: Optional[str], fallback: str = "6B7280") -> RGBColor:
    r, g, b = hex_to_tuple(hex_str, fallback)
    return RGBColor(r, g, b)


def normalize_theme(theme: ThemeConfig) -> ThemeConfig:
    """Defensive per-field fallback so one malformed hex code doesn't crash
    the whole deck or silently turn every color black (int('ZZ', 16) used
    to raise, uncaught, inside hex_to_rgb)."""
    defaults = {
        "bg_color": "0F172A",
        "text_color": "E2E8F0",
        "accent_1": "38BDF8",
        "accent_2": "F472B6",
        "panel_color": "1E293B",
    }
    for field, fallback in defaults.items():
        val = getattr(theme, field, None)
        cleaned = (val or "").strip().lstrip('#')
        if not _HEX_RE.match(cleaned):
            logger.warning(f"Theme field '{field}'='{val}' is not a valid hex color, defaulting to #{fallback}")
            setattr(theme, field, fallback)
        else:
            setattr(theme, field, cleaned)
    if not theme.font_name:
        theme.font_name = "Arial"
    return theme


# ---------------------------------------------------------------------------
# Image sourcing
#
# fetch_and_add_image() below is the drop-in replacement for the old
# network-only fetcher. It always returns True because it always has an
# image to place -- real photo if any source works, an on-theme generated
# placeholder if none do. Nothing upstream needs to special-case failure
# anymore.
# ---------------------------------------------------------------------------

_image_cache: Dict[str, bytes] = {}


def _clean_keyword(keyword: str) -> str:
    clean = re.sub(r'[^a-zA-Z0-9 ]', ' ', keyword or '').strip()
    return clean or "abstract background"


def _looks_like_image(content: bytes, content_type: str = "") -> bool:
    if content_type and not content_type.lower().startswith("image/"):
        return False
    if not content or len(content) < 800:
        return False
    try:
        Image.open(BytesIO(content)).verify()
        return True
    except Exception:
        return False


def _http_get_image(url: Optional[str], headers: Optional[dict] = None) -> Optional[bytes]:
    if not url:
        return None
    try:
        resp = requests.get(url, headers=headers or HTTP_HEADERS, timeout=IMAGE_TIMEOUT, allow_redirects=True)
        content_type = resp.headers.get("Content-Type", "")
        if resp.status_code == 200 and _looks_like_image(resp.content, content_type):
            return resp.content
    except requests.RequestException as e:
        logger.info(f"Image GET failed for {url}: {e}")
    except Exception as e:
        logger.info(f"Unexpected error fetching {url}: {e}")
    return None


def _fetch_pexels(keyword: str) -> Optional[bytes]:
    if not PEXELS_API_KEY:
        return None
    try:
        resp = requests.get(
            "https://api.pexels.com/v1/search",
            params={"query": keyword, "per_page": 6, "orientation": "landscape"},
            headers={"Authorization": PEXELS_API_KEY},
            timeout=IMAGE_TIMEOUT,
        )
        resp.raise_for_status()
        photos = resp.json().get("photos", [])
        if not photos:
            return None
        photo = random.choice(photos)
        src = (photo.get("src") or {}).get("large") or (photo.get("src") or {}).get("original")
        return _http_get_image(src)
    except Exception as e:
        logger.info(f"Pexels fetch failed for '{keyword}': {e}")
        return None


def _fetch_pixabay(keyword: str) -> Optional[bytes]:
    if not PIXABAY_API_KEY:
        return None
    try:
        resp = requests.get(
            "https://pixabay.com/api/",
            params={
                "key": PIXABAY_API_KEY, "q": keyword, "image_type": "photo",
                "orientation": "horizontal", "safesearch": "true", "per_page": 6,
            },
            timeout=IMAGE_TIMEOUT,
        )
        resp.raise_for_status()
        hits = resp.json().get("hits", [])
        if not hits:
            return None
        hit = random.choice(hits)
        src = hit.get("largeImageURL") or hit.get("webformatURL")
        return _http_get_image(src)
    except Exception as e:
        logger.info(f"Pixabay fetch failed for '{keyword}': {e}")
        return None


def _fetch_unsplash_api(keyword: str) -> Optional[bytes]:
    if not UNSPLASH_ACCESS_KEY:
        return None
    try:
        resp = requests.get(
            "https://api.unsplash.com/search/photos",
            params={"query": keyword, "per_page": 6, "orientation": "landscape"},
            headers={"Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}"},
            timeout=IMAGE_TIMEOUT,
        )
        resp.raise_for_status()
        results = resp.json().get("results", [])
        if not results:
            return None
        pick = random.choice(results)
        src = (pick.get("urls") or {}).get("regular")
        return _http_get_image(src)
    except Exception as e:
        logger.info(f"Unsplash API fetch failed for '{keyword}': {e}")
        return None


def _fetch_loremflickr(keyword: str) -> Optional[bytes]:
    # Real, keyword-tagged photos, no API key required.
    url = f"https://loremflickr.com/900/650/{requests.utils.quote(keyword)}"
    return _http_get_image(url)


def _fetch_picsum(keyword: str) -> Optional[bytes]:
    # NOTE: picsum has no keyword search. The "seed" only makes the pick
    # deterministic for a given string, not relevant to the topic. This is
    # kept as a last resort so a slide gets *some* photo, never as a
    # stand-in for an actual keyword match.
    url = f"https://picsum.photos/seed/{requests.utils.quote(keyword)}/900/650"
    return _http_get_image(url)


# source.unsplash.com intentionally removed: Unsplash deprecated it in
# 2021 and it has returned nothing but errors since. Keeping it in the
# chain was a guaranteed dead attempt on every single image fetch.
IMAGE_SOURCES = [_fetch_pexels, _fetch_pixabay, _fetch_unsplash_api, _fetch_loremflickr, _fetch_picsum]


def generate_placeholder_image(keyword: str, theme: ThemeConfig, width_px: int, height_px: int) -> bytes:
    """Guaranteed-to-succeed local visual so a slide is never left with a
    blank hole when every remote photo source fails (no outbound network,
    rate limits, everything down at once). Renders a soft on-theme abstract
    gradient with the deck's own accent colors instead of a gray box or a
    broken-image icon.
    """
    width_px = max(int(width_px), 80)
    height_px = max(int(height_px), 60)

    bg = hex_to_tuple(theme.panel_color or theme.bg_color)
    c1 = hex_to_tuple(theme.accent_1)
    c2 = hex_to_tuple(theme.accent_2)

    img = Image.new("RGB", (width_px, height_px), bg)
    draw = ImageDraw.Draw(img, "RGBA")

    # Diagonal soft wash from bg toward accent_1
    step = max(width_px // 200, 2)
    for x in range(0, width_px, step):
        t = x / max(width_px, 1)
        r = int(bg[0] * (1 - t * 0.4) + c1[0] * t * 0.4)
        g = int(bg[1] * (1 - t * 0.4) + c1[1] * t * 0.4)
        b = int(bg[2] * (1 - t * 0.4) + c1[2] * t * 0.4)
        draw.rectangle([x, 0, x + step, height_px], fill=(r, g, b))

    # A handful of soft translucent blobs in the accent colors, seeded by
    # the keyword so the same topic always gets the same look.
    rng = random.Random(sum(ord(c) for c in (keyword or "abstract")))
    short_side = min(width_px, height_px)
    for i in range(5):
        radius = rng.randint(int(short_side * 0.16), int(short_side * 0.40))
        cx = rng.randint(0, width_px)
        cy = rng.randint(0, height_px)
        color = c1 if i % 2 == 0 else c2
        draw.ellipse([cx - radius, cy - radius, cx + radius, cy + radius], fill=(color[0], color[1], color[2], 55))

    blur_radius = max(short_side // 40, 1)
    img = img.filter(ImageFilter.GaussianBlur(radius=blur_radius))

    out = BytesIO()
    img.convert("RGB").save(out, format="JPEG", quality=88)
    return out.getvalue()


def get_image_bytes(keyword: str, theme: ThemeConfig, target_px: Tuple[int, int]) -> bytes:
    """Always returns usable JPEG/PNG bytes. Tries real photo sources in
    order, and falls back to a generated on-theme placeholder so a slide is
    never left without a visual -- even with zero network access."""
    clean_kw = _clean_keyword(keyword)
    if clean_kw in _image_cache:
        return _image_cache[clean_kw]

    content = None
    for source_fn in IMAGE_SOURCES:
        for _ in range(IMAGE_MAX_RETRIES):
            content = source_fn(clean_kw)
            if content:
                break
        if content:
            logger.info(f"Image for '{clean_kw}' resolved via {source_fn.__name__}")
            break

    if not content:
        logger.warning(f"All remote image sources failed for '{clean_kw}'; generating local placeholder")
        content = generate_placeholder_image(clean_kw, theme, target_px[0], target_px[1])

    if len(_image_cache) > IMAGE_CACHE_MAX:
        _image_cache.clear()
    _image_cache[clean_kw] = content
    return content


def prefetch_slide_images(deck: PresentationDeck, theme: ThemeConfig) -> None:
    """Warms the image cache for every slide concurrently, so build_pptx
    doesn't block on a chain of sequential network round-trips -- one slow
    or hanging source used to slow down (or, with a short enough timeout,
    silently skip) every slide behind it."""
    keywords = []
    seen = set()
    for slide_content in deck.slides:
        visual = slide_content.visual
        if visual and visual.visual_type == "image" and visual.search_keyword:
            kw = _clean_keyword(visual.search_keyword)
            if kw not in seen:
                seen.add(kw)
                keywords.append(kw)

    if not keywords:
        return

    with ThreadPoolExecutor(max_workers=min(IMAGE_MAX_WORKERS, len(keywords))) as pool:
        futures = {pool.submit(get_image_bytes, kw, theme, (1200, 800)): kw for kw in keywords}
        for future in futures:
            try:
                future.result()
            except Exception as e:
                logger.warning(f"Prefetch failed for '{futures[future]}': {e}")


def place_image(slide, image_bytes: bytes, left_in: float, top_in: float, width_in: float, height_in: float) -> bool:
    """Center-crops the image to the target box's aspect ratio before
    inserting it, instead of stretching it to fit -- the old code passed
    add_picture a fixed width/height regardless of the source photo's
    native aspect ratio, which squashed or stretched almost every photo
    that wasn't already the right shape."""
    try:
        target_w = max(int(width_in * IMAGE_DPI), 10)
        target_h = max(int(height_in * IMAGE_DPI), 10)
        img = Image.open(BytesIO(image_bytes)).convert("RGB")
        fitted = ImageOps.fit(img, (target_w, target_h), method=Image.Resampling.LANCZOS)
        out = BytesIO()
        fitted.save(out, format="JPEG", quality=90)
        out.seek(0)
        slide.shapes.add_picture(out, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
        return True
    except Exception as e:
        logger.warning(f"Failed to place image: {e}")
        return False


def fetch_and_add_image(slide, keyword: str, theme: ThemeConfig, left_in: float, top_in: float,
                         width_in: float, height_in: float) -> bool:
    """Signature note: now takes plain inch floats (not Inches()-wrapped
    values) and an extra `theme` argument, since the guaranteed local
    fallback needs the theme's colors. Every render_slide call site below
    has been updated accordingly."""
    if not keyword:
        keyword = "abstract pattern"
    content = get_image_bytes(keyword, theme, (int(width_in * IMAGE_DPI), int(height_in * IMAGE_DPI)))
    return place_image(slide, content, left_in, top_in, width_in, height_in)


# ---------------------------------------------------------------------------
# Text-writing helpers shared across layouts
# ---------------------------------------------------------------------------

def setup_text_frame(tf):
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.18)


def write_bullet_list(tf, items: List[str], theme: ThemeConfig, font_size: int,
                       color_hex: Optional[str] = None, space_after: int = 10, prefix: str = "• "):
    """Writes into the text frame's existing first paragraph for item 0
    instead of always calling add_paragraph(), which used to leave a blank
    leading line in every bulleted box (two_column, hero, and the generic
    bullets fallback all had this)."""
    color_hex = color_hex or theme.text_color
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{prefix}{item}" if prefix else item
        p.font.name = theme.font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = hex_to_rgb(color_hex)
        p.space_after = Pt(space_after)


def write_title_desc_list(tf, items: List[DataItem], theme: ThemeConfig, box_w: float, box_h: float,
                           title_max_pt: int = 15, title_min_pt: int = 11,
                           desc_max_pt: int = 13, desc_min_pt: int = 10):
    if not items:
        return
    per_item_h = box_h / len(items)
    for i, item in enumerate(items):
        title = truncate_chars(item.title, 40)
        desc = truncate_chars(item.description, 140)
        title_size = fit_font_size(title, box_w, per_item_h * 0.4, max_pt=title_max_pt, min_pt=title_min_pt, bold=True)
        desc_size = fit_font_size(desc, box_w, per_item_h * 0.6, max_pt=desc_max_pt, min_pt=desc_min_pt)

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.name = theme.font_name
        p.font.size = Pt(title_size)
        p.font.color.rgb = hex_to_rgb(theme.accent_1)
        p.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = theme.font_name
        p2.font.size = Pt(desc_size)
        p2.font.color.rgb = hex_to_rgb(theme.text_color)
        p2.space_after = Pt(10)


def add_header(slide, title_text: str, subtitle_text: str, theme: ThemeConfig):
    header_w = 11.7
    title_text = truncate_chars(title_text, 70)
    title_size = fit_font_size(title_text, header_w, 0.8, max_pt=26, min_pt=18, bold=True, h_margin_in=0.1, v_margin_in=0.1)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(header_w), Inches(0.8))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = theme.font_name
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(theme.accent_1)

    if subtitle_text:
        subtitle_text = truncate_chars(subtitle_text, 110)
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(header_w), Inches(0.45))
        tf2 = txBox2.text_frame
        tf2.vertical_anchor = MSO_ANCHOR.TOP
        tf2.word_wrap = True
        tf2.auto_size = MSO_AUTO_SIZE.NONE
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.name = theme.font_name
        p2.font.size = Pt(13)
        p2.font.color.rgb = hex_to_rgb(theme.text_color)


def add_slide_number(slide, index: int, total: int, theme: ThemeConfig):
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.08), Inches(0.65), Inches(0.32))
    tf = txBox.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = f"{index}/{total}"
    p.font.name = theme.font_name
    p.font.size = Pt(9)
    p.font.color.rgb = hex_to_rgb(theme.text_color)
    p.alignment = PP_ALIGN.RIGHT


def pick_stat_value_and_label(item: DataItem) -> Tuple[str, str]:
    """The schema's generic title/description fields don't map cleanly onto
    a dashboard's 'big number + label'. Given how those fields are
    described ('title' = short, 'description' = a full sentence), the model
    reliably puts the short stat in `title` and the explanatory sentence in
    `description` -- confirmed against a real generated deck, where a
    30pt-bold card ended up displaying a full sentence while the actual
    stat ("65%+") was rendered as the small label underneath it. Rather
    than depend on the model matching the renderer's assumption, pick
    whichever field actually looks like a short stat.
    """
    t, d = (item.title or "").strip(), (item.description or "").strip()

    def is_stat_like(s: str) -> bool:
        bare = s.rstrip(".")
        return len(bare) <= 12 and bool(_STAT_RE.search(bare))

    t_stat, d_stat = is_stat_like(t), is_stat_like(d)
    if t_stat and not d_stat:
        return t, d
    if d_stat and not t_stat:
        return d, t
    # Ambiguous (both or neither look like a stat): the shorter string
    # reads better as the headline number either way.
    return (t, d) if len(t) <= len(d) else (d, t)


# ---------------------------------------------------------------------------
# Rendering
# ---------------------------------------------------------------------------

def add_background(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def render_slide(slide, slide_content: SlideContent, theme: ThemeConfig, prs: Presentation):
    add_background(slide, hex_to_rgb(theme.bg_color))
    layout = slide_content.layout
    has_image_kw = bool(slide_content.visual and slide_content.visual.search_keyword)

    if layout in ["title", "hero"]:
        is_pure_title = layout == "title" and not (slide_content.visual.visual_type == "image" and has_image_kw)

        if is_pure_title:
            # Centered Title Layout
            title_text = truncate_chars(slide_content.title, 80)
            title_size = fit_font_size(title_text, 11.33, 1.8, max_pt=40, min_pt=28, bold=True, h_margin_in=0.2, v_margin_in=0.15)

            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(1.8))
            tf = txBox.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            p.text = title_text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = theme.font_name
            p.font.size = Pt(title_size)
            p.font.bold = True
            p.font.color.rgb = hex_to_rgb(theme.accent_1)

            if slide_content.subtitle:
                subtitle_text = truncate_chars(slide_content.subtitle, 120)
                sub_size = fit_font_size(subtitle_text, 11.33, 1.2, max_pt=20, min_pt=14, h_margin_in=0.2, v_margin_in=0.15)
                txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.33), Inches(1.2))
                tf2 = txBox2.text_frame
                tf2.vertical_anchor = MSO_ANCHOR.TOP
                tf2.word_wrap = True
                tf2.auto_size = MSO_AUTO_SIZE.NONE
                p2 = tf2.paragraphs[0]
                p2.text = subtitle_text
                p2.alignment = PP_ALIGN.CENTER
                p2.font.name = theme.font_name
                p2.font.size = Pt(sub_size)
                p2.font.color.rgb = hex_to_rgb(theme.text_color)
        else:
            # Hero / Featured Image Layout
            add_header(slide, slide_content.title, slide_content.subtitle, theme)
            img_added = False
            if has_image_kw:
                img_added = fetch_and_add_image(slide, slide_content.visual.search_keyword, theme, 2.5, 1.8, 8.33, 4.8)

            if not img_added and slide_content.bullets:
                bullets = [truncate_chars(b, 130) for b in slide_content.bullets[:5]]
                box_w, box_h = 10.33, 4.5
                size = fit_font_size(" ".join(bullets), box_w, box_h, max_pt=18, min_pt=12)
                txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(box_w), Inches(box_h))
                tf = txBox.text_frame
                setup_text_frame(tf)
                tf.auto_size = MSO_AUTO_SIZE.NONE
                write_bullet_list(tf, bullets, theme, size, prefix="")

    elif layout == "two_column":
        add_header(slide, slide_content.title, slide_content.subtitle, theme)

        all_bullets = [truncate_chars(b, 110) for b in slide_content.bullets]
        left_bullets = all_bullets[:4]
        right_bullets = all_bullets[4:8]

        # Left column
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.0))
        tf = txBox.text_frame
        setup_text_frame(tf)
        tf.auto_size = MSO_AUTO_SIZE.NONE
        if left_bullets:
            size = fit_font_size(" ".join(left_bullets), 5.6, 5.0, max_pt=16, min_pt=11)
            write_bullet_list(tf, left_bullets, theme, size)
        elif slide_content.data_items:
            write_title_desc_list(tf, slide_content.data_items[:3], theme, 5.6, 5.0)

        # Right Column (Image or Secondary Content) -- no longer duplicates
        # bullets already shown on the left.
        img_added = False
        if has_image_kw:
            img_added = fetch_and_add_image(slide, slide_content.visual.search_keyword, theme, 6.8, 1.8, 5.7, 5.0)

        if not img_added:
            txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
            tf2 = txBox2.text_frame
            setup_text_frame(tf2)
            tf2.auto_size = MSO_AUTO_SIZE.NONE
            if right_bullets:
                size = fit_font_size(" ".join(right_bullets), 5.7, 5.0, max_pt=16, min_pt=11)
                write_bullet_list(tf2, right_bullets, theme, size)
            elif len(slide_content.data_items) > 3:
                write_title_desc_list(tf2, slide_content.data_items[3:6], theme, 5.7, 5.0)

    elif layout in ["cards", "infographic"]:
        add_header(slide, slide_content.title, slide_content.subtitle, theme)

        items = slide_content.data_items[:4] or [DataItem(title="Highlight", description="Key point about this topic.")]

        show_side_image = (slide_content.visual.visual_type == "image" and has_image_kw and len(items) <= 3)
        cards_total_width = 11.7
        if show_side_image:
            image_left, image_width = 8.2, 4.3
            img_ok = fetch_and_add_image(slide, slide_content.visual.search_keyword, theme, image_left, 1.8, image_width, 5.0)
            cards_total_width = 7.0 if img_ok else 11.7

        count = len(items)
        spacing = 0.35
        card_w = (cards_total_width - (count - 1) * spacing) / count
        card_h = 5.0

        for i, item in enumerate(items):
            left = 0.8 + i * (card_w + spacing)
            title = truncate_chars(item.title, 40)
            desc = truncate_chars(item.description, 160)

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(card_w), Inches(card_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(theme.panel_color)
            shape.line.color.rgb = hex_to_rgb(theme.accent_1)
            shape.line.width = Pt(1.5)

            tf = shape.text_frame
            setup_text_frame(tf)
            tf.auto_size = MSO_AUTO_SIZE.NONE

            title_size = fit_font_size(title, card_w, card_h * 0.22, max_pt=16, min_pt=12, bold=True)
            desc_size = fit_font_size(desc, card_w, card_h * 0.7, max_pt=13, min_pt=10)

            p = tf.paragraphs[0]
            p.text = title
            p.font.name = theme.font_name
            p.font.bold = True
            p.font.size = Pt(title_size)
            p.font.color.rgb = hex_to_rgb(theme.accent_1)
            p.space_after = Pt(6)

            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.name = theme.font_name
            p2.font.size = Pt(desc_size)
            p2.font.color.rgb = hex_to_rgb(theme.text_color)

    elif layout in ["architecture", "process_flow", "timeline", "roadmap"]:
        add_header(slide, slide_content.title, slide_content.subtitle, theme)

        items = slide_content.data_items[:5] or [DataItem(title="Step 1", description="Process description")]
        count = len(items)
        spacing = 0.4
        step_w = (11.7 - (count - 1) * spacing) / count
        step_h = 1.2

        for i, item in enumerate(items):
            left = 0.8 + i * (step_w + spacing)
            title = truncate_chars(item.title, 30)
            desc = truncate_chars(item.description, 130)

            shape_type = MSO_SHAPE.HEXAGON if layout == "architecture" else (
                MSO_SHAPE.OVAL if layout in ["timeline", "roadmap"] else MSO_SHAPE.RECTANGLE
            )
            shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(1.8), Inches(step_w), Inches(step_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(theme.accent_2)
            shape.line.color.rgb = hex_to_rgb(theme.accent_1)
            shape.line.width = Pt(1.5)

            tf = shape.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.05)

            # Hexagons/ovals taper at the edges, so the usable text width is
            # narrower than the shape's own bounding box -- otherwise a
            # longer title visually crosses the outline even though it
            # technically fits the rectangle python-pptx thinks it owns.
            width_factor = 0.70 if shape_type == MSO_SHAPE.HEXAGON else (0.80 if shape_type == MSO_SHAPE.OVAL else 0.90)
            title_size = fit_font_size(title, step_w * width_factor, step_h, max_pt=14, min_pt=10, bold=True, h_margin_in=0.15, v_margin_in=0.15)

            p = tf.paragraphs[0]
            p.text = title
            p.font.name = theme.font_name
            p.font.bold = True
            p.font.size = Pt(title_size)
            p.font.color.rgb = hex_to_rgb(theme.bg_color)
            p.alignment = PP_ALIGN.CENTER

            desc_h = 3.6
            desc_size = fit_font_size(desc, step_w, desc_h, max_pt=12, min_pt=10)
            txBox = slide.shapes.add_textbox(Inches(left), Inches(3.2), Inches(step_w), Inches(desc_h))
            tf2 = txBox.text_frame
            setup_text_frame(tf2)
            tf2.auto_size = MSO_AUTO_SIZE.NONE
            p2 = tf2.paragraphs[0]
            p2.text = desc
            p2.font.name = theme.font_name
            p2.font.size = Pt(desc_size)
            p2.font.color.rgb = hex_to_rgb(theme.text_color)
            p2.alignment = PP_ALIGN.CENTER

            if i < count - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + step_w + 0.08), Inches(2.25), Inches(spacing - 0.16), Inches(0.3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = hex_to_rgb(theme.accent_1)
                arrow.line.fill.background()

    elif layout == "dashboard":
        add_header(slide, slide_content.title, slide_content.subtitle, theme)

        items = slide_content.data_items[:4] or [DataItem(title="100+", description="Metric description")]

        for i, item in enumerate(items):
            col, row = i % 2, i // 2
            left = 0.8 + col * 6.0
            top = 1.8 + row * 2.6
            card_w, card_h = 5.7, 2.3

            # FIX: the old code assumed `description` was always the big
            # headline number and `title` the label -- backwards from how
            # the model actually fills these fields. This picks whichever
            # field is actually stat-shaped.
            value_text, label_text = pick_stat_value_and_label(item)
            value_text = truncate_chars(value_text, 30)
            label_text = truncate_chars(label_text, 90)

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(card_w), Inches(card_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(theme.panel_color)
            shape.line.color.rgb = hex_to_rgb(theme.accent_2)
            shape.line.width = Pt(1.5)

            tf = shape.text_frame
            setup_text_frame(tf)
            tf.auto_size = MSO_AUTO_SIZE.NONE

            value_size = fit_font_size(value_text, card_w, card_h * 0.55, max_pt=30, min_pt=16, bold=True)
            label_size = fit_font_size(label_text, card_w, card_h * 0.4, max_pt=13, min_pt=10)

            p = tf.paragraphs[0]
            p.text = value_text
            p.font.name = theme.font_name
            p.font.bold = True
            p.font.size = Pt(value_size)
            p.font.color.rgb = hex_to_rgb(theme.accent_1)
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(4)

            p2 = tf.add_paragraph()
            p2.text = label_text
            p2.font.name = theme.font_name
            p2.font.size = Pt(label_size)
            p2.font.color.rgb = hex_to_rgb(theme.text_color)
            p2.alignment = PP_ALIGN.CENTER

    elif layout == "comparison":
        add_header(slide, slide_content.title, slide_content.subtitle, theme)

        items = slide_content.data_items[:2] or [DataItem(title="Option A", description="Details A"), DataItem(title="Option B", description="Details B")]
        card_w, card_h = 5.6, 5.0

        for i, item in enumerate(items):
            left = 0.8 + i * 6.1
            title = truncate_chars(item.title, 30)
            desc = truncate_chars(item.description, 220)
            color = theme.accent_1 if i == 0 else theme.accent_2

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(card_w), Inches(card_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(theme.panel_color)
            shape.line.color.rgb = hex_to_rgb(color)
            shape.line.width = Pt(2.5)

            tf = shape.text_frame
            setup_text_frame(tf)
            tf.auto_size = MSO_AUTO_SIZE.NONE

            title_size = fit_font_size(title, card_w, card_h * 0.15, max_pt=18, min_pt=14, bold=True)
            desc_size = fit_font_size(desc, card_w, card_h * 0.8, max_pt=13, min_pt=10)

            p = tf.paragraphs[0]
            p.text = title
            p.font.name = theme.font_name
            p.font.bold = True
            p.font.size = Pt(title_size)
            p.font.color.rgb = hex_to_rgb(color)
            p.space_after = Pt(8)

            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.name = theme.font_name
            p2.font.size = Pt(desc_size)
            p2.font.color.rgb = hex_to_rgb(theme.text_color)

    else:  # bullets or generic fallback
        add_header(slide, slide_content.title, slide_content.subtitle, theme)

        show_side_image = (slide_content.visual.visual_type == "image" and has_image_kw)
        text_width = 5.6 if show_side_image else 11.7

        bullets = [truncate_chars(b, 140) for b in slide_content.bullets[:6]]
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(text_width), Inches(5.0))
        tf = txBox.text_frame
        setup_text_frame(tf)
        tf.auto_size = MSO_AUTO_SIZE.NONE
        if bullets:
            size = fit_font_size(" ".join(bullets), text_width, 5.0, max_pt=17, min_pt=11)
            write_bullet_list(tf, bullets, theme, size)

        if show_side_image:
            fetch_and_add_image(slide, slide_content.visual.search_keyword, theme, 6.8, 1.8, 5.7, 5.0)


def render_fallback_slide(slide, slide_content: SlideContent, theme: ThemeConfig):
    """Used when render_slide() throws for any reason. Deliberately simple
    (header + plain bullets, no shapes, no images) so it has as little
    surface area to fail on as possible -- the goal is that the rest of the
    deck still ships even if one slide's content or layout was pathological."""
    try:
        add_background(slide, hex_to_rgb(theme.bg_color))
        add_header(slide, slide_content.title or "Untitled Slide", slide_content.subtitle, theme)
        raw_items = slide_content.bullets or [d.description for d in slide_content.data_items]
        bullets = [truncate_chars(b, 140) for b in raw_items[:6] if b]
        if bullets:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
            tf = txBox.text_frame
            setup_text_frame(tf)
            size = fit_font_size(" ".join(bullets), 11.7, 5.0, max_pt=16, min_pt=11)
            write_bullet_list(tf, bullets, theme, size)
    except Exception as e:
        logger.error(f"Fallback slide rendering also failed, shipping a blank slide: {e}")


# ---------------------------------------------------------------------------
# Deck assembly
# ---------------------------------------------------------------------------

def normalize_deck(deck: PresentationDeck) -> PresentationDeck:
    """Defensive clean-up applied right after the LLM returns structured
    output, independent of whether the model actually followed the
    prompt's word limits on a given call. These ceilings are looser than
    the per-layout ones inside render_slide (which are tuned to each box);
    this pass just guarantees the renderer never sees something truly
    pathological, e.g. a 2,000-character 'description'."""
    for slide in deck.slides:
        slide.title = truncate_chars(slide.title, 90) or "Untitled"
        slide.subtitle = truncate_chars(slide.subtitle, 140)
        slide.bullets = [truncate_chars(b, 160) for b in (slide.bullets or [])][:8]
        for item in slide.data_items:
            item.title = truncate_chars(item.title, 60)
            item.description = truncate_chars(item.description, 220)
        if slide.visual and not slide.visual.search_keyword:
            slide.visual.search_keyword = (slide.title.split(" ")[0] if slide.title else "abstract")
    return deck


def build_pptx(deck: PresentationDeck) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    theme = normalize_theme(deck.theme)
    prefetch_slide_images(deck, theme)

    blank_layout = prs.slide_layouts[6]
    total = len(deck.slides)

    for idx, slide_content in enumerate(deck.slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        try:
            render_slide(slide, slide_content, theme, prs)
        except Exception as e:
            logger.error(
                f"Slide {idx} ('{slide_content.title}', layout='{slide_content.layout}') failed to "
                f"render: {e}. Falling back to a plain-text slide so the rest of the deck still ships."
            )
            render_fallback_slide(slide, slide_content, theme)

        add_slide_number(slide, idx, total, theme)

        if slide_content.speaker_notes:
            try:
                slide.notes_slide.notes_text_frame.text = slide_content.speaker_notes
            except Exception as e:
                logger.warning(f"Could not write speaker notes for slide {idx}: {e}")

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def safe_filename(title: str) -> str:
    name = re.sub(r"[^A-Za-z0-9_-]+", "_", title or "").strip("_")[:60].rstrip("_")
    return f"{name or 'OmniMind_Presentation'}.pptx"


def _invoke_with_retry(chain, messages, max_attempts: int = 3, base_delay: float = 1.5):
    """The old code made one LLM call and gave up on any exception,
    including transient ones (a rate limit, a network blip). Retries with
    a short linear backoff before surfacing the error."""
    last_err = None
    for attempt in range(1, max_attempts + 1):
        try:
            return chain.invoke(messages)
        except Exception as e:
            last_err = e
            logger.warning(f"LLM call attempt {attempt}/{max_attempts} failed: {e}")
            if attempt < max_attempts:
                time.sleep(base_delay * attempt)
    raise last_err


def ppt_agent(state):
    logger.info("===== PPT Agent =====")

    try:
        feedback = state.get("review_feedback", "")
        system_msg = SYSTEM_PROMPT
        if feedback:
            system_msg += f"\n\nCRITICAL FIXES REQUIRED BASED ON REVIEW FEEDBACK:\n{feedback}"

        user_prompt = state["query"]
        collected = state.get("collected_requirements", {})
        if collected:
            req_details = "\n".join([f"- {k}: {v}" for k, v in collected.items() if v])
            user_prompt = f"User Query: {user_prompt}\n\nCollected Requirements & Topic Details:\n{req_details}"

        chain = get_ppt_chain()
        deck: PresentationDeck = _invoke_with_retry(
            chain,
            [SystemMessage(content=system_msg), HumanMessage(content=user_prompt)],
        )

        deck = normalize_deck(deck)
        pptx_bytes = build_pptx(deck)
        encoded_pptx = base64.b64encode(pptx_bytes).decode("utf-8")

        return {
            "final_response": {
                "type": "ppt",
                "title": deck.title,
                "subtitle": "",
                "theme": deck.theme.style_name,
                "slide_count": len(deck.slides),
                "file_name": safe_filename(deck.title),
                "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "base64_data": encoded_pptx,
                "outline": [
                    {
                        "title": slide.title,
                        "slide_type": slide.layout,
                    }
                    for slide in deck.slides
                ],
            }
        }

    except Exception as e:
        logger.error(f"PPT Agent Error: {e}", exc_info=True)

        return {
            "final_response": {
                "type": "error",
                "message": f"Could not generate the PowerPoint deck: {str(e)}",
            }
        }